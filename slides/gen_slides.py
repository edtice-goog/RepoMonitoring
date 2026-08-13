"""Generate the RepoMonitoring vision deck (Black Duck SCA integration pitch).

Usage:  python gen_slides.py
Output: RepoMonitoring-Vision.pptx (same directory)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- palette
INK = RGBColor(0x21, 0x21, 0x2B)        # near-black text
GRAY = RGBColor(0x6B, 0x6B, 0x76)       # secondary text
PURPLE = RGBColor(0x58, 0x2C, 0x83)     # Black Duck-adjacent accent
PURPLE_SOFT = RGBColor(0xED, 0xE7, 0xF6)
RED = RGBColor(0xC0, 0x39, 0x2B)        # attacker / threat
RED_SOFT = RGBColor(0xFB, 0xE9, 0xE7)
GREEN = RGBColor(0x2E, 0x7D, 0x32)      # defended outcome
GREEN_SOFT = RGBColor(0xE8, 0xF5, 0xE9)
SLATE = RGBColor(0x45, 0x4A, 0x55)      # neutral boxes
SLATE_SOFT = RGBColor(0xEC, 0xEF, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def new_slide():
    return prs.slides.add_slide(BLANK)


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb


def para(tf, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         level=0, space_after=6, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def slide_header(slide, title, kicker=None):
    if kicker:
        tb = textbox(slide, Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.4))
        para(tb.text_frame, kicker.upper(), size=13, color=PURPLE, bold=True, first=True)
        ty = Inches(0.68)
    else:
        ty = Inches(0.45)
    tb = textbox(slide, Inches(0.55), ty, Inches(12.2), Inches(0.9))
    para(tb.text_frame, title, size=32, color=INK, bold=True, first=True)


def footer(slide, n):
    tb = textbox(slide, Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.35))
    para(tb.text_frame, f"RepoMonitoring — vision concept for Black Duck SCA integration   |   {n}",
         size=10, color=GRAY, first=True)


def box(slide, x, y, w, h, title, body=None, fill=PURPLE_SOFT, edge=PURPLE,
        title_color=None, body_color=GRAY, title_size=14, body_size=11.5,
        align=PP_ALIGN.CENTER):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = edge
    shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    para(tf, title, size=title_size, color=title_color or edge, bold=True,
         align=align, first=True, space_after=2)
    if body:
        para(tf, body, size=body_size, color=body_color, align=align, space_after=0)
    return shp


def arrow(slide, x1, y1, x2, y2, color=SLATE, width=2.25, dash=None):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.shadow.inherit = False
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    if dash:
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.insert(0, d)
    return conn


def mid(shape, side):
    """Edge midpoint of a shape: 'l', 'r', 't', 'b'."""
    if side == "l":
        return shape.left, shape.top + shape.height // 2
    if side == "r":
        return shape.left + shape.width, shape.top + shape.height // 2
    if side == "t":
        return shape.left + shape.width // 2, shape.top
    return shape.left + shape.width // 2, shape.top + shape.height


def connect(slide, a, b, side_a="r", side_b="l", **kw):
    x1, y1 = mid(a, side_a)
    x2, y2 = mid(b, side_b)
    return arrow(slide, x1, y1, x2, y2, **kw)


# ================================================================ slide 1 — title
s = new_slide()
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE; bar.line.fill.background()
bar.shadow.inherit = False

tb = textbox(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.6))
para(tb.text_frame, "Closing the Patch Gap for Embedded Builds", size=44, bold=True,
     color=INK, first=True)
para(tb.text_frame, "Upstream security-fix monitoring as a Black Duck SCA component",
     size=22, color=PURPLE)

tb = textbox(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(1.6))
para(tb.text_frame,
     "Alerting Yocto / Buildroot device makers to likely security fixes the moment they "
     "land upstream — before a CVE exists, and only when the fix touches code they actually shipped.",
     size=16, color=GRAY, first=True)
para(tb.text_frame, "Vision deck — concept + working prototype", size=13, color=GRAY)

# ================================================================ slide 2 — the threat
s = new_slide()
slide_header(s, "Exploits are now derived from the patches that fix them",
             kicker="The problem")

tb = textbox(s, Inches(0.55), Inches(1.55), Inches(12.2), Inches(0.7))
para(tb.text_frame,
     "Attackers diff public fix commits and weaponize them in days. Vendors who build entire "
     "operating systems from source are exposed the longest: rebuild, regression, and fleet "
     "update cycles are measured in weeks.", size=15, color=GRAY, first=True)

# timeline
y = Inches(3.1); bh = Inches(1.15); bw = Inches(2.5); gap = Inches(0.45)
xs = [Inches(0.7) + i * (bw + gap) for i in range(4)]
t1 = box(s, xs[0], y, bw, bh, "Fix merged upstream",
         "Often silent — no advisory,\nno CVE, just a commit", fill=SLATE_SOFT, edge=SLATE)
t2 = box(s, xs[1], y, bw, bh, "Exploit derived",
         "Attackers diff the patch;\nN-day weaponized in days", fill=RED_SOFT, edge=RED)
t3 = box(s, xs[2], y, bw, bh, "CVE published",
         "Weeks later — the first moment\nfeed-driven tools can alert", fill=SLATE_SOFT, edge=SLATE)
t4 = box(s, xs[3], y, bw, bh, "Vendor ships fix",
         "Rebuild + regression + fleet\nupdate: weeks to months", fill=SLATE_SOFT, edge=SLATE)
for a, b in [(t1, t2), (t2, t3), (t3, t4)]:
    connect(s, a, b)

# exposure window brace
ex = box(s, xs[1], Inches(4.75), xs[3] + bw - xs[1], Inches(0.85),
         "Exposure window", "Devices are exploitable in the field for this entire span",
         fill=RED_SOFT, edge=RED)

tb = textbox(s, Inches(0.55), Inches(5.95), Inches(12.2), Inches(0.8))
para(tb.text_frame,
     "The defensive opportunity: the same public commit that arms the attacker can warn the "
     "vendor — if someone is watching the right repos with the right context.",
     size=15, color=PURPLE, bold=True, first=True)
footer(s, 2)

# ================================================================ slide 3 — why SCA misses it
s = new_slide()
slide_header(s, "CVE-driven alerting can't cover this window", kicker="The gap in SCA today")

rows = [
    ("Alerts begin at CVE publication",
     "The exposure window opens at the fix commit — weeks earlier. By publication, "
     "N-day exploitation is already underway."),
    ("Component-level granularity",
     "A kernel CVE alerts every kernel user. A vendor whose kernel config never builds the "
     "affected driver gets a false alarm — for hardware the device doesn't have."),
    ("Alert fatigue defeats the tool",
     "Embedded teams triaging vulnerability feeds learn to ignore them. "
     "Precision is the product, not just coverage."),
]
y = Inches(1.7)
for title, body in rows:
    b1 = box(s, Inches(0.7), y, Inches(3.6), Inches(1.35), title, fill=PURPLE_SOFT,
             edge=PURPLE, align=PP_ALIGN.LEFT, title_size=15)
    tb = textbox(s, Inches(4.6), y + Inches(0.12), Inches(8.1), Inches(1.2))
    para(tb.text_frame, body, size=14.5, color=INK, first=True)
    y += Inches(1.62)

tb = textbox(s, Inches(0.7), Inches(6.45), Inches(12.0), Inches(0.55))
para(tb.text_frame,
     "Needed: a pre-CVE signal scoped to what each customer's build actually compiled.",
     size=16, color=PURPLE, bold=True, first=True)
footer(s, 3)

# ================================================================ slide 4 — the insight
s = new_slide()
slide_header(s, "Monitor the build, not the feed", kicker="The insight")

c1 = box(s, Inches(0.7), Inches(1.7), Inches(5.9), Inches(2.3),
         "RECALL — never miss a repo",
         "SBOM (SCA) merged with build-observation component matches "
         "produces the complete list of upstream repos that feed the build. "
         "Two independent sources: neither alone is complete.",
         fill=PURPLE_SOFT, edge=PURPLE, title_size=16, body_size=14, body_color=INK)
c2 = box(s, Inches(6.85), Inches(1.7), Inches(5.9), Inches(2.3),
         "PRECISION — never false-alarm",
         "Every upstream change is cross-referenced against the files actually "
         "compiled into the image (sources and headers). Changes to code the "
         "device never shipped are suppressed — with the reason logged.",
         fill=GREEN_SOFT, edge=GREEN, title_size=16, body_size=14, body_color=INK)

ex = box(s, Inches(0.7), Inches(4.35), Inches(12.05), Inches(1.5),
         "Worked example",
         "A security fix lands in a kernel Wi-Fi driver. Repo-level mapping says “exposed.” "
         "The compiled-file index knows the customer's kernel config never built that driver — "
         "no alert, no noise, confidence preserved for the alerts that do fire.",
         fill=SLATE_SOFT, edge=SLATE, title_size=15, body_size=14, body_color=INK,
         align=PP_ALIGN.LEFT)

tb = textbox(s, Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.6))
para(tb.text_frame,
     "Precision comes from observing the real build — not trusting declared dependencies.",
     size=16, color=PURPLE, bold=True, first=True)
footer(s, 4)

# ================================================================ slide 5 — BD already has the data
s = new_slide()
slide_header(s, "Black Duck already collects the hard part", kicker="Why Black Duck SCA")

tb = textbox(s, Inches(0.7), Inches(1.65), Inches(12.0), Inches(2.7))
for txt, first in [
    ("The expensive input is the compiled-file inventory — and the BD C/C++ tool "
     "(Coverity Build Capture) already produces it on every analyzed build:", True),
]:
    para(tb.text_frame, txt, size=16, color=INK, first=first)
for b in [
    "Wraps the customer's build, observing every compiler and linker invocation",
    "Records all compiled sources, included headers, and linked objects",
    "Matches files to components via package manager, file contents, and binary analysis",
    "Already submits the complete file list to Black Duck SCA as part of normal operation",
]:
    para(tb.text_frame, "•  " + b, size=15, color=INK, level=1, space_after=8)

ex = box(s, Inches(0.7), Inches(4.6), Inches(12.05), Inches(1.6),
         "Integration cost is connecting existing data to a new consumer",
         "No new collection agents, no new customer workflow. The build the customer already "
         "analyzes with BD/CPP yields both the SBOM (recall) and the compiled-file index "
         "(precision). Header capture even covers fixes that land in .h files.",
         fill=PURPLE_SOFT, edge=PURPLE, title_size=16, body_size=14, body_color=INK)
footer(s, 5)

# ================================================================ slide 6 — vision architecture
s = new_slide()
slide_header(s, "Vision: a monitoring service inside the BD SCA deployment",
             kicker="Architecture")

# lane labels
for label, x, w, col in [
    ("CUSTOMER BUILD ENV", Inches(0.55), Inches(3.3), SLATE),
    ("BLACK DUCK SCA (HUB)", Inches(4.25), Inches(3.7), PURPLE),
    ("MONITORING SERVICE  —  new deployment container", Inches(8.35), Inches(4.45), GREEN),
]:
    tb = textbox(s, x, Inches(1.55), w, Inches(0.35))
    para(tb.text_frame, label, size=12, color=col, bold=True, align=PP_ALIGN.CENTER, first=True)

# lane backgrounds
for x, w, col in [(Inches(0.55), Inches(3.3), SLATE), (Inches(4.25), Inches(3.7), PURPLE),
                  (Inches(8.35), Inches(4.45), GREEN)]:
    lane = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.95), w, Inches(4.55))
    lane.fill.solid(); lane.fill.fore_color.rgb = WHITE
    lane.line.color.rgb = col; lane.line.width = Pt(1)
    lane.shadow.inherit = False
    d = lane.line._get_or_add_ln()
    dash = d.makeelement(qn("a:prstDash"), {"val": "dash"})
    d.append(dash)

b_build = box(s, Inches(0.8), Inches(2.5), Inches(2.8), Inches(1.5),
              "BD/CPP wraps the build",
              "Coverity Build Capture:\nsources + headers + objects",
              fill=SLATE_SOFT, edge=SLATE)
b_hub = box(s, Inches(4.5), Inches(2.5), Inches(3.2), Inches(1.5),
            "Project BOM + matched file list",
            "Existing BD/CPP submission —\nno new collection",
            fill=PURPLE_SOFT, edge=PURPLE)
b_notify = box(s, Inches(4.5), Inches(4.55), Inches(3.2), Inches(1.5),
               "New build analyzed",
               "Repo list computed; push\nnotification to monitoring service",
               fill=PURPLE_SOFT, edge=PURPLE)
b_watch = box(s, Inches(8.6), Inches(2.35), Inches(3.95), Inches(1.25),
              "Watch manifest per project",
              "Repos, pinned versions,\ncompiled-file index",
              fill=GREEN_SOFT, edge=GREEN)
b_repos = box(s, Inches(8.6), Inches(3.8), Inches(3.95), Inches(1.25),
              "Repo watchers",
              "Webhooks (modern) / polling (legacy)\n— transport invisible to the user",
              fill=GREEN_SOFT, edge=GREEN)
b_flow = box(s, Inches(8.6), Inches(5.25), Inches(3.95), Inches(1.05),
             "Change detected → workflow",
             "Filter → triage → notification",
             fill=GREEN_SOFT, edge=GREEN)

connect(s, b_build, b_hub)
connect(s, b_hub, b_notify, side_a="b", side_b="t")
connect(s, b_notify, b_watch)
connect(s, b_watch, b_repos, side_a="b", side_b="t")
connect(s, b_repos, b_flow, side_a="b", side_b="t")

tb = textbox(s, Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.5))
para(tb.text_frame,
     "Ships and scales like the rest of BD SCA: one more container in the existing deployment.",
     size=14, color=PURPLE, bold=True, first=True)
footer(s, 6)

# ================================================================ slide 7 — workflow
s = new_slide()
slide_header(s, "From upstream commit to customer notification", kicker="Detection workflow")

y = Inches(2.0); bh = Inches(1.25); bw = Inches(2.35); gap = Inches(0.42)
xs = [Inches(0.6) + i * (bw + gap) for i in range(4)]
w1 = box(s, xs[0], y, bw, bh, "Commit detected",
         "In a watched repo,\non a watched branch", fill=SLATE_SOFT, edge=SLATE)
w2 = box(s, xs[1], y, bw, bh, "Relevance filter",
         "Changed files × compiled-\nfile index (incl. headers)", fill=PURPLE_SOFT, edge=PURPLE)
w3 = box(s, xs[2], y, bw, bh, "Version check",
         "Does the change apply to\nthe version that was built?", fill=PURPLE_SOFT, edge=PURPLE)
w4 = box(s, xs[3], y, bw, bh, "Triage",
         "Human first;\nLLM-assisted → agentic", fill=PURPLE_SOFT, edge=PURPLE)
connect(s, w1, w2); connect(s, w2, w3); connect(s, w3, w4)

sup = box(s, xs[1], Inches(3.85), bw, Inches(0.95), "Suppressed + logged",
          "Code not in this build\n(the kernel-driver case)", fill=SLATE_SOFT, edge=SLATE)
connect(s, w2, sup, side_a="b", side_b="t", color=GRAY, width=1.75)

# verdicts
vy = Inches(3.85); vh = Inches(0.95); vw = Inches(2.35)
v1 = box(s, Inches(11.55) - 2 * (vw + Inches(0.25)), vy, vw, vh, "Not meaningful",
         "Rationale logged", fill=SLATE_SOFT, edge=SLATE)
v2 = box(s, Inches(11.55) - (vw + Inches(0.25)), vy, vw, vh, "Needs human review",
         "Queued with rationale", fill=PURPLE_SOFT, edge=PURPLE)
v3 = box(s, Inches(11.55), vy, Inches(1.78), vh, "Response required",
         "Urgency + rationale", fill=RED_SOFT, edge=RED)
for v in (v1, v2, v3):
    connect(s, w4, v, side_a="b", side_b="t", color=GRAY, width=1.5)

out = box(s, Inches(7.2), Inches(5.35), Inches(5.55), Inches(1.15),
          "Notification → candidate rebuild",
          "Evidence bundle: commit, matched files, applicability, reasoning. "
          "Customer starts rebuild + regression before any zero-day lands.",
          fill=GREEN_SOFT, edge=GREEN, body_size=12.5)
connect(s, v3, out, side_a="b", side_b="t", color=GREEN)
connect(s, v2, out, side_a="b", side_b="t", color=GREEN, width=1.5)

tb = textbox(s, Inches(0.6), Inches(5.5), Inches(6.2), Inches(1.0))
para(tb.text_frame,
     "Every verdict carries reasoning — auditable decisions, not a black box. "
     "Scope is urgency triage only: the tool never derives exploits.",
     size=13.5, color=GRAY, first=True)
footer(s, 7)

# ================================================================ slide 8 — lifecycle
s = new_slide()
slide_header(s, "Always in sync with the latest analyzed build", kicker="Lifecycle")

steps = [
    ("Customer re-runs BD/CPP on a new build", "The workflow they already have — nothing new to operate"),
    ("Hub recomputes the project's repo list", "From the fresh BOM + matched file list"),
    ("Push notification to the monitoring service", "Watch lists adjusted: repos added, dropped, re-pinned"),
    ("Compiled-file index refreshed", "Precision filter always reflects what the device actually runs"),
]
y = Inches(1.8)
for i, (t, b) in enumerate(steps, 1):
    num = box(s, Inches(0.7), y, Inches(0.55), Inches(0.55), str(i),
              fill=PURPLE, edge=PURPLE, title_color=WHITE, title_size=18)
    tb = textbox(s, Inches(1.5), y - Inches(0.05), Inches(11.2), Inches(0.9))
    para(tb.text_frame, t, size=17, color=INK, bold=True, first=True, space_after=2)
    para(tb.text_frame, b, size=13.5, color=GRAY)
    y += Inches(1.05)

ex = box(s, Inches(0.7), Inches(6.0), Inches(12.05), Inches(0.85),
         "Zero user configuration",
         "Webhook vs. polling per repo, branch selection, version pinning — all invisible. "
         "The user's only experience is the notification.",
         fill=GREEN_SOFT, edge=GREEN, title_size=14, body_size=12.5, body_color=INK)
footer(s, 8)

# ================================================================ slide 9 — why BD wins
s = new_slide()
slide_header(s, "Why this belongs in Black Duck SCA", kicker="Product case")

cards = [
    ("Differentiated signal", "Pre-CVE security-fix alerting scoped to the customer's actual "
     "build. No CVE-feed product can match the timing; no repo-watcher can match the precision."),
    ("Near-zero collection cost", "BD/CPP already captures and submits the compiled-file "
     "inventory. The feature monetizes data the product already holds."),
    ("Fits the deployment model", "One additional container in the existing BD SCA deployment. "
     "No new agents at the customer."),
    ("The right segment", "Yocto / Buildroot device makers carry the longest patch gap and the "
     "highest exposure — and are already BD/CPP's core users."),
]
positions = [(Inches(0.7), Inches(1.8)), (Inches(6.85), Inches(1.8)),
             (Inches(0.7), Inches(4.25), ), (Inches(6.85), Inches(4.25))]
for (t, b), (x, y) in zip(cards, positions):
    box(s, x, y, Inches(5.9), Inches(2.15), t, b, fill=PURPLE_SOFT, edge=PURPLE,
        title_size=17, body_size=13.5, body_color=INK, align=PP_ALIGN.LEFT)
footer(s, 9)

# ================================================================ slide 10 — roadmap + demo
s = new_slide()
slide_header(s, "Rollout path — and what exists today", kicker="Roadmap")

ph = [
    ("Phase 1", "Human triage", "Filtered, evidence-rich review queue. The relevance filter "
     "does the noise reduction; analysts make the call."),
    ("Phase 2", "LLM-assisted", "Claude classifies each surviving commit with a verdict and "
     "written rationale; humans review the queue it produces."),
    ("Phase 3", "Agentic", "End-to-end triage with human oversight on response_required "
     "verdicts only. Notification triggers the rebuild pipeline."),
]
x = Inches(0.7)
for tag, t, b in ph:
    c = box(s, x, Inches(1.8), Inches(3.9), Inches(2.3), f"{tag} — {t}", b,
            fill=PURPLE_SOFT, edge=PURPLE, title_size=15.5, body_size=13,
            body_color=INK, align=PP_ALIGN.LEFT)
    x += Inches(4.12)

demo = box(s, Inches(0.7), Inches(4.5), Inches(12.05), Inches(1.5),
           "Working prototype (today)",
           "End-to-end pipeline on sample data: provisioning from an SBOM + build-capture file "
           "list, replayable change injection, relevance filtering, three-verdict triage with "
           "rationale, and the alert evidence bundle. Fully deterministic — built for live demo.",
           fill=GREEN_SOFT, edge=GREEN, title_size=16, body_size=14, body_color=INK)

tb = textbox(s, Inches(0.7), Inches(6.25), Inches(12.0), Inches(0.6))
para(tb.text_frame,
     "Ask: sponsor a proof-of-concept integration against a real BD/CPP-analyzed project.",
     size=17, color=PURPLE, bold=True, first=True)
footer(s, 10)

# ---------------------------------------------------------------- save
out = Path(__file__).parent / "RepoMonitoring-Vision.pptx"
prs.save(out)
print(f"Wrote {out}")
