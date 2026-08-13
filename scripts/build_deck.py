#!/usr/bin/env python3
"""Generate the Patch Gap executive deck (.pptx) from structured content.

Content mirrors docs/patch-gap-exec-deck.md. Re-run after edits:
    python scripts/build_deck.py
Output: docs/patch-gap-exec-deck.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- theme ---------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x38, 0x64)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)   # restrained red for the punch numbers
GRAY = RGBColor(0x7F, 0x7F, 0x7F)
DARK = RGBColor(0x26, 0x26, 0x26)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def _set(run, size, color=DARK, bold=False, italic=False, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def add_title_bar(slide, title):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    _set(r, 32, NAVY, bold=True)
    # accent underline
    line = slide.shapes.add_shape(1, Inches(0.62), Inches(1.32), Inches(2.2), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    return box


def add_source(slide, source):
    if not source:
        return
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(12.1), Inches(0.4))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Source: " + source
    _set(r, 11, GRAY, italic=True)


def add_notes(slide, notes):
    if not notes:
        return
    slide.notes_slide.notes_text_frame.text = notes


def add_bullets(slide, bullets, top=1.7, height=5.0, base=20):
    """bullets: list of (text, level, kind) where kind in {bullet, plain, lead}."""
    box = slide.shapes.add_textbox(Inches(0.7), Inches(top), Inches(12.0), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    first = True
    for item in bullets:
        text, level, kind = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        if kind == "lead":
            r = p.add_run(); r.text = text
            _set(r, base + 2, NAVY, bold=True)
        elif kind == "plain":
            r = p.add_run(); r.text = text
            _set(r, base, DARK)
        else:  # bullet
            glyph = "    – " if level else "•  "
            r = p.add_run(); r.text = glyph + text
            _set(r, base - (2 if level else 0), DARK)
    return box


# --- slide builders ------------------------------------------------------
def slide_title(title, subtitle, notes):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)

    box = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.0))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    _set(r, 54, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
    _set(r2, 24, RGBColor(0xCF, 0xD6, 0xE4), italic=True)
    add_notes(s, notes)


def slide_content(title, bullets, source, notes, base=20, top=1.7):
    s = prs.slides.add_slide(BLANK)
    add_title_bar(s, title)
    add_bullets(s, bullets, top=top, base=base)
    add_source(s, source)
    add_notes(s, notes)


def slide_tte(title, rows, caption, bullet, source, notes):
    s = prs.slides.add_slide(BLANK)
    add_title_bar(s, title)
    cols, nrows = 2, len(rows)
    tbl_shape = s.shapes.add_table(nrows, cols, Inches(0.7), Inches(1.8),
                                   Inches(6.2), Inches(0.6 * nrows))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(3.6)
    tbl.columns[1].width = Inches(2.6)
    for ri, (c0, c1, emph) in enumerate(rows):
        for ci, txt in enumerate((c0, c1)):
            cell = tbl.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = cell.text_frame.paragraphs[0]
            run = para.add_run(); run.text = txt
            if ri == 0:
                _set(run, 18, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            else:
                _set(run, 18, ACCENT if emph else DARK, bold=emph)
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT if ri % 2 else RGBColor(0xFF, 0xFF, 0xFF)
    # caption + forward-looking bullet on the right
    box = s.shapes.add_textbox(Inches(7.3), Inches(1.9), Inches(5.4), Inches(4.5))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = caption
    _set(r, 18, NAVY, italic=True, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(18)
    r2 = p2.add_run(); r2.text = bullet
    _set(r2, 16, DARK)
    add_source(s, source)
    add_notes(s, notes)


# --- content -------------------------------------------------------------
slide_title(
    "The Patch Gap",
    "Defending the window before a CVE exists",
    "Thesis up front: the most dangerous moment isn't when a vulnerability is announced — "
    "it's the quiet window between the fix becoming visible and our product being rebuilt. "
    "That window is measurable, shrinking, and widest for embedded. Everything that follows is sourced.")

slide_content(
    "The Exposure Window",
    [("A security fix becomes public knowledge the moment it's committed — not when the CVE is announced. Attackers race us through that gap.", 0, "lead"),
     ("The fix itself is the disclosure. The code change shows exactly what was broken.", 0, "bullet"),
     ("We don't ship when upstream fixes a bug — we have to rebuild and re-release firmware.", 0, "bullet"),
     ("The difference between those two timestamps is our exposure window.", 0, "bullet")],
    "Google Project Zero (2020).",
    "The whole thesis. If the room remembers one slide, it's this one. The rest is evidence that attackers exploit this window, that it's collapsing, and that it's widest for embedded.")

slide_content(
    "How the attack works",
    [("\"Patch diffing\" — reverse-engineering the fix to find the flaw:", 0, "lead"),
     ("Vendor publishes an update (or an OSS project merges a commit).", 1, "bullet"),
     ("Attacker compares \"before\" vs. \"after\" — the change points straight at the vulnerability.", 1, "bullet"),
     ("Attacker builds a working exploit from that diff.", 1, "bullet"),
     ("Anyone still running the old version is now a target.", 1, "bullet"),
     ("Routine enough to have a nickname: \"Patch Tuesday, Exploit Wednesday.\"", 0, "bullet"),
     ("Self-evident: a defect found, reported, and fixed comes with a map — the commit is the map. Re-finding it is a low bar.", 0, "lead")],
    "IBM X-Force (2023): CVE-2023-21768 weaponized in ~1 day. Google Project Zero (2020): ~3-week upper bound.",
    "Concrete but non-technical. IBM X-Force diffed a public Windows patch with off-the-shelf tools and had a working exploit in about a day — with no prior experience in that component. The technique is a standard skill, not exotic.",
    base=18)

slide_tte(
    "Time-to-exploit is collapsing",
    [("Period", "Avg. days to exploit", False),
     ("2018–2019", "63 days", False),
     ("2020–2021", "44 days", False),
     ("2021–2022", "32 days", False),
     ("2023", "5 days", True)],
    "\"This is less than a sixth of the previously observed TTE.\" — Mandiant",
    "A floor, not a ceiling: the data ends in 2023. Building exploits keeps getting cheaper and faster — an accelerant on an already-downward line. The point stands on its own; no extra study required.",
    "Mandiant / Google Cloud, 2023 Time-to-Exploit Trends (Oct 2024).",
    "Headline chart. Five years ago defenders had two months; in 2023 the average was five days. Caveat in your pocket: the 5-day average excludes 15 outliers (~47 days with them) and blends zero-day and n-day timing; the direction is corroborated independently.")

slide_content(
    "N-day: exploitation after a patch exists",
    [("Of 138 vulnerabilities exploited in 2023, 41 (30%) were N-days — first exploited after a patch existed.", 0, "bullet"),
     ("Of those N-days: 12% within 1 day, 29% within 1 week, 56% within 1 month of the patch.", 0, "bullet"),
     ("This is the patch-gap window, measured directly.", 0, "lead")],
    "Mandiant / Google Cloud, 2023 Time-to-Exploit report (Oct 2024).",
    "Pairs with the prior slide. N-day is the precise thing we defend against: attacks on people who had a fix available but hadn't shipped it. For an embedded vendor on a months-long rebuild cycle, \"within a month\" means the attack lands before we've shipped. Raw counts if pressed: 5 / 12 / 23 of 41.")

slide_content(
    "Silent fixes: patches precede disclosure",
    [("Open-source fixes are routinely public in the code before any CVE or advisory exists:", 0, "lead"),
     ("~70% of security patches are committed before public disclosure.", 0, "bullet"),
     ("Vulnerability databases (Snyk, NVD) lag a median of 25 days behind the fix.", 0, "bullet"),
     ("~38% of security releases ship with no note that they fix a security issue.", 0, "bullet"),
     ("Linux kernel: fixes detectable in the public repo 2–179 days before disclosure.", 0, "bullet")],
    "Li & Paxson, ACM CCS 2017; Imtiaz et al., IEEE TSE 2022; Ramsauer et al., ACM CCSW 2020.",
    "Why \"just subscribe to a CVE feed\" is not enough — by the time the CVE exists, the attacker has had the diff for weeks. Peer-reviewed, not vendor marketing. The Linux number matters most: it's the exact ecosystem under embedded firmware.")

slide_content(
    "Embedded carries the widest exposure",
    [("Forescout teardown of real router / edge firmware (2024):", 0, "lead"),
     ("Average open-source component is 5 years, 6 months old — and 4 years, 4 months behind latest.", 0, "bullet"),
     ("161 known vulnerabilities baked into the average firmware image (24 rated critical).", 0, "bullet"),
     ("~20 exploitable N-day kernel vulnerabilities per image.", 0, "bullet"),
     ("4 of 5 images analyzed ran OpenWrt (embedded Linux) — our ecosystem.", 0, "bullet")],
    "Forescout, \"Rough Around the Edges\" OT/IoT router firmware study (Aug 2024).",
    "Makes it personal. Embedded products carry huge, old dependency trees; every one of those 161 known vulnerabilities had a public fix that never reached the shipped image. Long rebuild cycles + large dependency trees = the widest exposure window of anyone.")

slide_content(
    "Attackers concentrate on edge devices",
    [("Over 60% of enterprise zero-day exploitation in 2024 hit security & network devices (20 of 33) — Ivanti, Palo Alto, Cisco ASA, Fortinet, VMware.", 0, "bullet"),
     ("~half of enterprise zero-days in 2025, same pattern.", 0, "bullet"),
     ("Why: edge / embedded devices can't run endpoint detection (EDR) — Mandiant calls this \"a blind spot for defenders… an ideal attack surface.\"", 0, "bullet")],
    "Mandiant / Google Cloud, 2024 Zero-Day Trends & 2025 Zero-Day Review.",
    "Two reinforcing facts: attackers are concentrating on the device class we build, and that class is where defenders are blindest because you can't install agents on a router or camera. The defensive answer to an on-device blind spot is to watch the supply chain upstream — which is what our tool does.")

slide_content(
    "Scale and business impact",
    [("+180% (≈3×) year-over-year surge in breaches that start with vulnerability exploitation — Verizon DBIR 2024.", 0, "bullet"),
     ("Organizations take 55 days to remediate half of critical vulns after a patch ships — DBIR 2024. (Attackers need 5.)", 0, "bullet"),
     ("CISA Known-Exploited-Vulnerabilities catalog: 1,484 entries, +245 (+20%) in 2025 — its largest jump in three years.", 0, "bullet"),
     ("Average breach cost $4.88M, +10% YoY (largest jump since the pandemic) — IBM 2024.", 0, "bullet"),
     ("PAN-OS CVE-2024-3400: public PoC one day after disclosure; a dozen-plus groups exploiting it within two weeks.", 0, "bullet")],
    "Verizon DBIR 2024; SecurityWeek/CISA KEV 2025; IBM Cost of a Data Breach 2024; Mandiant M-Trends 2025.",
    "The \"so what, in dollars and trend lines\" slide. 55-days-vs-5-days is the knockout: defenders patch in 55 days, attackers exploit in 5. PAN-OS is the named, recent case. Note: DBIR/KEV/IBM figures are primary-sourced but were not run through the adversarial verification pass.",
    base=18)

slide_content(
    "A recognized field: tools, standards, research",
    [("The direct answer to \"where's the evidence?\"", 0, "lead"),
     ("CISA maintains a government catalog of actively-exploited vulns with mandated remediation deadlines (KEV / BOD).", 0, "bullet"),
     ("Google OSV, OpenSSF cve-bin-tool, Timesys Vigiles — production tools to find vulnerable components in builds (incl. Buildroot & Yocto).", 0, "bullet"),
     ("\"Security Patch Detection\" is a named academic discipline — a 2024 ACM Computing Surveys paper reviews 127 studies (2014–2023).", 0, "bullet")],
    "CISA KEV; google/osv.dev; ossf/cve-bin-tool; TimesysGit/vigiles-buildroot; Lin et al., ACM Computing Surveys 2024.",
    "This slide exists because of the pushback. The threat model isn't our invention — governments legislate around it, Google and the Linux Foundation ship tools against it, and there's a decade of peer-reviewed literature. What's novel in our work is applying it as continuous monitoring for our specific supply chain.")

slide_content(
    "Where our work fits",
    [("RepoMonitoring closes the gap on the one axis we control: defender reaction time.", 0, "lead"),
     ("Watch the upstream repos our products depend on — including commits, before any CVE.", 0, "bullet"),
     ("Triage each change: is this a security fix or a normal bug fix?", 0, "bullet"),
     ("Alert our build teams to rebuild during the gap — not 25 days later when the CVE finally posts.", 0, "bullet"),
     ("We can't make attackers slower. We can stop being the slowest one in the room.", 0, "lead")],
    "Internal — RepoMonitoring project scope.",
    "Tie everything together. Slide 6 says the signal exists in the code before the CVE; slide 10 says detecting it is solved enough to build on; this slide operationalizes that for our dependency tree. Defensive, triage-only — we flag what to rebuild; we never derive exploits.")

slide_content(
    "The ask",
    [("Recognize patch-gap exposure as a tracked risk for our embedded products — not just \"patch when the CVE lands.\"", 0, "bullet"),
     ("Resource continuous upstream monitoring + commit triage for our core dependency set.", 0, "bullet"),
     ("Target metric: shrink our fix-visible → firmware-shipped window — the one number this deck is about.", 0, "bullet")],
    None,
    "Keep the ask concrete and small. We're not boiling the ocean; we're measuring and shrinking one window for a defined dependency set. Close by returning to the trend: the industry went from 63 days to 5. How many days are we, and are we moving the right way?")

slide_content(
    "Appendix A — Sources",
    [("Primary threat-intelligence", 0, "lead"),
     ("Mandiant / Google Cloud — 2023 Time-to-Exploit Trends (Oct 2024); 2024 Zero-Day Trends; 2025 Zero-Day Review; M-Trends 2025.", 1, "bullet"),
     ("Google Project Zero (2020); IBM X-Force (2023).", 1, "bullet"),
     ("Embedded / scale / cost", 0, "lead"),
     ("Forescout, Rough Around the Edges (2024); Verizon DBIR 2024; CISA KEV via SecurityWeek (2025); IBM Cost of a Data Breach 2024.", 1, "bullet"),
     ("Peer-reviewed academic", 0, "lead"),
     ("Li & Paxson, ACM CCS 2017; Imtiaz et al., IEEE TSE 2022; Ramsauer et al., ACM CCSW 2020; Lin et al., ACM Computing Surveys 2024.", 1, "bullet"),
     ("Tooling", 0, "lead"),
     ("google/osv.dev; ossf/cve-bin-tool; TimesysGit/vigiles-buildroot; CISA KEV catalog.", 1, "bullet")],
    None,
    "Full citation list with org + year. Every figure in the deck traces here.",
    base=18, top=1.6)

slide_content(
    "Appendix B — Anticipated objections",
    [("\"The 5-day number sounds cherry-picked.\" → Excludes 15 outliers; with them ~47 days, still down from 63. Corroborated independently. The trend is the point.", 0, "bullet"),
     ("\"Those are vendor marketing reports.\" → Methodology & samples are published; the mechanism is independently established in peer-reviewed venues (CCS, TSE, CCSW).", 0, "bullet"),
     ("\"Silent-fix percentages are general OSS, not our embedded stack.\" → Correct, and we say so. Embedded exposure is proven separately by Forescout + Mandiant edge data. They meet at the Linux kernel.", 0, "bullet"),
     ("\"Show me one concrete case.\" → PAN-OS CVE-2024-3400: public PoC one day after disclosure, a dozen-plus groups within two weeks.", 0, "bullet"),
     ("Do not use: \"Linux fixes precede CVE by ~100 days (Kroah-Hartman)\" — this claim was refuted in our fact-check.", 0, "bullet")],
    None,
    "Insurance against a repeat of the last presentation. Every likely objection has a one-line sourced rebuttal. The do-not-use entry keeps a tempting-but-unverified stat out of your mouth.",
    base=17)

out = Path(__file__).resolve().parent.parent / "docs" / "patch-gap-exec-deck.pptx"
prs.save(str(out))
print(f"Wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
